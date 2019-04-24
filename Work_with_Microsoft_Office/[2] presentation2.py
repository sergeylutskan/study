from pptx import Presentation


prs = Presentation()
layout = prs.slide_layouts[0]
for i in range(5):
    prs.slides.add_slide(layout)


ft = prs.slides[0].shapes[0].text_frame
ft2 = prs.slides[0].shapes[1].text_frame
p = ft.paragraphs[0]
p2 = ft2.paragraphs[0]
run = p.add_run()
run2 = p2.add_run()
run.text = "choice(seq) method of Random instance"
run2.text = "Choose a random element from a non-empty sequence."
font = run.font
font2 = run2.font
font.name = 'Courier New'
font2.name = 'Courier New'

ft = prs.slides[1].shapes[0].text_frame
ft2 = prs.slides[1].shapes[1].text_frame
p = ft.paragraphs[0]
p2 = ft2.paragraphs[0]
run = p.add_run()
run2 = p2.add_run()
run.text = "random(...) method of Random instance"
run2.text = "random() -> x in the interval [0, 1)."
font = run.font
font2 = run2.font
font.name = 'Courier New'
font2.name = 'Courier New'

ft = prs.slides[2].shapes[0].text_frame
ft2 = prs.slides[2].shapes[1].text_frame
p = ft.paragraphs[0]
p2 = ft2.paragraphs[0]
run = p.add_run()
run2 = p2.add_run()
run.text = "uniform(a, b) method of Random instance"
run2.text = "Get a random number in the range [a, b) or [a, b] depending on rounding."
font = run.font
font2 = run2.font
font.name = 'Courier New'
font2.name = 'Courier New'

ft = prs.slides[3].shapes[0].text_frame
ft2 = prs.slides[3].shapes[1].text_frame
p = ft.paragraphs[0]
p2 = ft2.paragraphs[0]
run = p.add_run()
run2 = p2.add_run()
run.text = "setstate(state) method of Random instance"
run2.text = "Restore internal state from object returned by getstate()."
font = run.font
font2 = run2.font
font.name = 'Courier New'
font2.name = 'Courier New'

ft = prs.slides[4].shapes[0].text_frame
ft2 = prs.slides[4].shapes[1].text_frame
p = ft.paragraphs[0]
p2 = ft2.paragraphs[0]
run = p.add_run()
run2 = p2.add_run()
run.text = "randint(a, b) method of Random instance"
run2.text = "Return random integer in range [a, b], including both end points."
font = run.font
font2 = run2.font
font.name = 'Courier New'
font2.name = 'Courier New'

prs.save('test.pptx')
