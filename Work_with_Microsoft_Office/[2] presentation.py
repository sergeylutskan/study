from pptx import Presentation
 
prs = Presentation()
 
layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

run = title.add_run()
font = run.font
font.name = 'Courier New'
#title.text = "choice(seq) method of Random instance"

subtitle.text = '''
Choose a random element from a non-empty sequence.
'''

slide = prs.slides.add_slide(layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
 
title.text = "random(...) method of Random instance"

subtitle.text = '''
random() -> x in the interval [0, 1).
'''
slide = prs.slides.add_slide(layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
 
title.text = "uniform(a, b) method of Random instance"

subtitle.text = '''
Get a random number in the range [a, b) or [a, b] depending on rounding.
'''
slide = prs.slides.add_slide(layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
 
title.text = "setstate(state) method of Random instance"

subtitle.text = '''
Restore internal state from object returned by getstate().
'''
slide = prs.slides.add_slide(layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
 
title.text = "randint(a, b) method of Random instance"

subtitle.text = '''
Return random integer in range [a, b], including both end points.
'''
prs.save('test.pptx')
